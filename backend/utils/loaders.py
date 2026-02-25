import os
import io
import pytesseract
from PIL import Image
from pptx import Presentation
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    UnstructuredPowerPointLoader
)

# --- CRITICAL RAILWAY FIX ---
# Remove the Windows-specific Tesseract path. 
# On Linux (Railway), the binary will be found automatically in the system PATH.
# ----------------------------

def load_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return PyPDFLoader(file_path).load()
    elif ext == ".docx":
        return Docx2txtLoader(file_path).load()
    elif ext in [".ppt", ".pptx"]:
        return load_pptx_with_pages(file_path)
    elif ext in [".png", ".jpg", ".jpeg"]:
        return load_image_with_ocr(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def load_pptx_with_pages(file_path):
    prs = Presentation(file_path)
    documents = []
    
    for i, slide in enumerate(prs.slides):
        combined_text = []
        has_text = False
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                combined_text.append(shape.text)
                has_text = True # Mark that this slide already has text
            
            # THE SPEED FIX: Only run OCR if the slide has very little or no text
            elif shape.shape_type == 13 and not has_text: 
                try:
                    image_bytes = shape.image.blob
                    image_stream = io.BytesIO(image_bytes)
                    image_doc = load_image_with_ocr_from_stream(image_stream)
                    if image_doc[0].page_content.strip():
                        combined_text.append(f"[Image Content: {image_doc[0].page_content}]")
                except Exception as e:
                    print(f"DEBUG: Failed to OCR image on slide {i+1}: {e}")

        content = "\n".join(combined_text)

        if content.strip():
            documents.append(Document(
                page_content=content,
                metadata={"source": file_path, "page": i + 1}
            ))
    return documents

def load_image_with_ocr_from_stream(image_stream):
    try:
        image = Image.open(image_stream).convert("L")
        image = image.point(lambda x: 0 if x < 140 else 255, '1')
        text = pytesseract.image_to_string(image, config="--psm 6")
        return [Document(page_content=text)]
    except:
        return [Document(page_content="")]
        
def load_image_with_ocr(image_path):
    if os.path.getsize(image_path) > 3_000_000:
        return [Document(page_content="", metadata={"source": image_path, "type": "image", "ocr_confidence": 0.0})]

    try:
        # Pre-processing for better OCR
        image = Image.open(image_path).convert("L")
        image = image.point(lambda x: 0 if x < 140 else 255, '1')

        # Tesseract will be called via the Linux system binary
        text = pytesseract.image_to_string(image, config="--psm 6")

        text = text.encode("utf-8", errors="ignore").decode("utf-8")
        text = "".join(c for c in text if c.isprintable())
    except Exception as e:
        text = ""
        print(f"OCR failed for {image_path}: {e}")

    confidence = estimate_ocr_confidence(text)
    return [Document(page_content=text, metadata={"source": image_path, "type": "image", "ocr_confidence": confidence})]

def estimate_ocr_confidence(text):
    if not text or len(text.strip()) < 20:
        return 0.2
    elif len(text) < 100:
        return 0.5
    else:
        return 0.8



